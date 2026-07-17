"""
Annual Report and Investor Presentation Parser.

Parses company annual reports, quarterly reports, and investor presentations.
Extracts structured financial data, KPIs, strategic initiatives, and outlook.
"""

import asyncio
import logging
import re
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any, Optional

from ..financial_documents.parser import parse_financial_pdf, ParserResult
from ..chunking.section_splitter import create_chunker

logger = logging.getLogger(__name__)


@dataclass
class AnnualReportData:
    """Structured data extracted from an annual report."""
    company_name: str
    ticker: Optional[str] = None
    fiscal_year: Optional[int] = None
    filing_date: Optional[date] = None
    cik: Optional[str] = None
    
    # Business overview
    business_description: str = ""
    products_services: list[str] = field(default_factory=list)
    markets: list[str] = field(default_factory=list)
    competitive_position: str = ""
    
    # Financial highlights
    revenue: Optional[float] = None
    net_income: Optional[float] = None
    eps: Optional[float] = None
    gross_margin: Optional[float] = None
    operating_margin: Optional[float] = None
    free_cash_flow: Optional[float] = None
    
    # Key metrics
    key_metrics: dict[str, Any] = field(default_factory=dict)
    
    # Business segments
    segments: list[dict] = field(default_factory=list)
    
    # Risk factors
    risk_factors: list[str] = field(default_factory=list)
    
    # MD&A highlights
    mda_highlights: dict[str, str] = field(default_factory=dict)
    
    # Capital allocation
    dividends: Optional[float] = None
    buybacks: Optional[float] = None
    capex: Optional[float] = None
    
    # Guidance
    guidance: dict[str, Any] = field(default_factory=dict)
    
    # ESG
    esg_highlights: dict[str, Any] = field(default_factory=dict)
    
    # Source
    source_file: Optional[str] = None
    pages_processed: int = 0
    extraction_confidence: float = 0.0


@dataclass
class QuarterlyReportData:
    """Structured data extracted from a quarterly report (10-Q)."""
    company_name: str
    ticker: Optional[str] = None
    quarter: Optional[str] = None  # Q1, Q2, Q3, Q4
    fiscal_year: Optional[int] = None
    filing_date: Optional[date] = None
    
    # Financial results
    revenue: Optional[float] = None
    revenue_yoy_change: Optional[float] = None
    net_income: Optional[float] = None
    eps: Optional[float] = None
    eps_yoy_change: Optional[float] = None
    gross_margin: Optional[float] = None
    operating_margin: Optional[float] = None
    
    # Guidance
    guidance: dict[str, Any] = field(default_factory=dict)
    
    # Key highlights
    highlights: list[str] = field(default_factory=list)
    
    # Segment performance
    segment_performance: list[dict] = field(default_factory=list)
    
    # MD&A highlights
    mda_summary: str = ""


@dataclass
class InvestorPresentationData:
    """Structured data extracted from investor presentations."""
    company_name: str
    ticker: Optional[str] = None
    presentation_date: Optional[date] = None
    presentation_title: str = ""
    
    # Slides data
    slides: list[dict] = field(default_factory=list)
    
    # Key sections
    key_highlights: list[str] = field(default_factory=list)
    financial_highlights: dict[str, Any] = field(default_factory=dict)
    strategic_initiatives: list[str] = field(default_factory=list)
    growth_drivers: list[str] = field(default_factory=list)
    risks: list[str] = field(default_factory=list)
    
    # Financial metrics shown
    key_metrics: dict[str, Any] = field(default_factory=dict)
    
    # Guidance
    guidance: dict[str, Any] = field(default_factory=dict)
    
    # Capital allocation
    capital_allocation: dict[str, Any] = field(default_factory=dict)
    
    # ESG
    esg_highlights: dict[str, Any] = field(default_factory=dict)
    
    # Source
    source_file: Optional[str] = None
    total_slides: int = 0


class AnnualReportParser:
    """
    Parse company annual reports (10-K, Annual Reports).
    
    Extracts:
    - Business overview
    - Financial highlights
    - Segment information
    - Risk factors
    - MD&A key points
    - Capital allocation
    - Guidance
    """
    
    def __init__(self):
        self.chunker = create_chunker()
    
    async def parse(self, file_path: Path) -> AnnualReportData:
        """Parse annual report PDF or HTML."""
        # Parse document
        result = await parse_financial_pdf(file_path)
        
        if not result.success:
            raise ValueError(f"Failed to parse {file_path}: {result.errors}")
        
        # Combine text
        full_text = "\n\n".join(page.text for page in result.pages)
        
        # Create base data
        data = AnnualReportData(
            source_file=str(file_path),
            pages_processed=len(result.pages)
        )
        
        # Extract metadata from first pages
        await self._extract_metadata(full_text, data, result)
        
        # Extract business overview
        await self._extract_business_overview(full_text, data)
        
        # Extract financial highlights
        await self._extract_financial_highlights(full_text, data)
        
        # Extract segments
        await self._extract_segments(full_text, data)
        
        # Extract MD&A highlights
        await self._extract_mda_highlights(full_text, data)
        
        # Extract risk factors
        await self._extract_risk_factors(full_text, data)
        
        # Extract capital allocation
        await self._extract_capital_allocation(full_text, data)
        
        # Extract guidance
        await self._extract_guidance(full_text, data)
        
        # Calculate confidence
        data.extraction_confidence = self._calculate_confidence(data)
        
        return data
    
    async def _extract_metadata(self, text: str, data: AnnualReportData, result: ParserResult):
        """Extract company metadata from document."""
        # Try to find company name
        lines = text.split('\n')[:50]
        for line in lines:
            if any(keyword in line.lower() for keyword in ['corporation', 'inc.', 'inc', 'company', 'ltd', 'llc']):
                data.company_name = line.strip()
                break
        
        # Try to find ticker
        ticker_match = re.search(r'\(([A-Z]{1,5})\)', text[:5000])
        if ticker_match:
            data.ticker = ticker_match.group(1)
        
        # Try to find CIK
        cik_match = re.search(r'cik[:\s]*(\d{10})', text, re.IGNORECASE)
        if cik_match:
            data.cik = cik_match.group(1)
        
        # Find fiscal year
        year_matches = re.findall(r'(?:fiscal\s+year|year\s+ended)\s+(?:20\d{2})', text, re.IGNORECASE)
        if year_matches:
            year_match = re.search(r'(20\d{2})', year_matches[0])
            if year_match:
                data.fiscal_year = int(year_match.group(1))
        
        # Filing date
        date_match = re.search(r'filed\s+(?:on\s+)?(\w+\s+\d{1,2},?\s+20\d{2})', text, re.IGNORECASE)
        if date_match:
            try:
                data.filing_date = datetime.strptime(date_match.group(1), "%B %d, %Y").date()
            except:
                pass
    
    async def _extract_business_overview(self, text: str, data: AnnualReportData):
        """Extract business description and overview."""
        # Find business section
        patterns = [
            r'(?i)(?:item\s+1\.?\s*business|business\s+overview|our\s+business)',
            r'(?i)(?:who\s+we\s+are|company\s+overview)',
        ]
        
        business_text = self._extract_section(text, patterns, max_chars=10000)
        
        if business_text:
            data.business_description = business_text[:5000]
            
            # Extract products/services
            product_matches = re.findall(
                r'(?:products?|services?|solutions?|offerings?)[:\\s]+([^\n]{20,200})',
                business_text, re.IGNORECASE
            )
            data.products_services = [p.strip() for p in product_matches[:10]]
            
            # Extract markets
            market_matches = re.findall(
                r'(?:markets?|customers?|industries?|sectors?)[:\\s]+([^\n]{20,200})',
                business_text, re.IGNORECASE
            )
            data.markets = [m.strip() for m in market_matches[:10]]
    
    async def _extract_financial_highlights(self, text: str, data: AnnualReportData):
        """Extract key financial highlights."""
        # Revenue
        rev_match = re.search(
            r'(?:total\s+)?revenue\s+(?:of|was|is|amounted to)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if rev_match:
            value = float(rev_match.group(1).replace(',', ''))
            if 'billion' in rev_match.group(0).lower() or 'b ' in rev_match.group(0).lower():
                value *= 1e9
            elif 'million' in rev_match.group(0).lower() or 'm ' in rev_match.group(0).lower():
                value *= 1e6
            data.revenue = value
        
        # Net income
        ni_match = re.search(
            r'net income\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if ni_match:
            value = float(ni_match.group(1).replace(',', ''))
            if 'billion' in ni_match.group(0).lower():
                value *= 1e9
            elif 'million' in ni_match.group(0).lower():
                value *= 1e6
            data.net_income = value
        
        # EPS
        eps_match = re.search(
            r'(?:eps|earnings per share)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)',
            text, re.IGNORECASE
        )
        if eps_match:
            data.eps = float(eps_match.group(1).replace(',', ''))
        
        # Margins
        gm_match = re.search(r'gross margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%', text, re.IGNORECASE)
        if gm_match:
            data.gross_margin = float(gm_match.group(1))
        
        om_match = re.search(r'operating margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%', text, re.IGNORECASE)
        if om_match:
            data.operating_margin = float(om_match.group(1))
        
        # Free cash flow
        fcf_match = re.search(
            r'free cash flow\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if fcf_match:
            value = float(fcf_match.group(1).replace(',', ''))
            if 'billion' in fcf_match.group(0).lower():
                value *= 1e9
            elif 'million' in fcf_match.group(0).lower():
                value *= 1e6
            data.free_cash_flow = value
    
    async def _extract_segments(self, text: str, data: AnnualReportData):
        """Extract business segment information."""
        # Look for segment reporting section
        segment_patterns = [
            r'(?i)(?:segment\s+information|reportable\s+segments?|operating\s+segments?)',
            r'(?i)item\s+2\.?\s*properties',  # Sometimes segments in Item 2
        ]
        
        segment_text = self._extract_section(text, segment_patterns, max_chars=15000)
        
        if segment_text:
            # Look for segment names and metrics
            # Pattern: Segment name followed by revenue/profit
            segment_pattern = re.compile(
                r'([A-Z][A-Za-z\s]{3,30})\s*(?:segment|division)?[:\\s]+.*?(?:revenue|sales|income).*?[\$£€]?\s*([\d,\.]+)',
                re.IGNORECASE | re.DOTALL
            )
            
            for match in re.finditer(segment_pattern, segment_text):
                segment_name = match.group(1).strip()
                revenue_str = match.group(2).replace(',', '')
                try:
                    revenue = float(revenue_str)
                    data.segments.append({
                        'name': segment_name,
                        'revenue': revenue,
                        'type': 'operating_segment'
                    })
                except:
                    pass
    
    async def _extract_mda_highlights(self, text: str, data: AnnualReportData):
        """Extract key MD&A highlights."""
        # Find MD&A section
        mda_patterns = [
            r'(?i)(?:item\s+7\.?\s*management\'?s\s+discussion\s+and\s+analysis|md&a|management\'s\s+discussion)',
        ]
        
        mda_text = self._extract_section(text, mda_patterns, max_chars=20000)
        
        if mda_text:
            # Extract key topics
            highlights = {}
            
            # Liquidity and capital resources
            liquidity_text = self._extract_subsection(mda_text, [
                r'(?i)(?:liquidity\s+and\s+capital\s+resources|capital\s+resources\s+and\s+liquidity)'
            ])
            if liquidity_text:
                highlights['liquidity'] = liquidity_text[:2000]
            
            # Results of operations
            results_text = self._extract_subsection(mda_text, [
                r'(?i)(?:results\s+of\s+operations|operating\s+results)'
            ])
            if results_text:
                highlights['results_of_operations'] = results_text[:2000]
            
            # Critical accounting estimates
            critical_text = self._extract_subsection(mda_text, [
                r'(?i)(?:critical\s+accounting\s+(?:estimates|policies))'
            ])
            if critical_text:
                highlights['critical_accounting'] = critical_text[:2000]
            
            # Contractual obligations
            contract_text = self._extract_subsection(mda_text, [
                r'(?i)(?:contractual\s+obligations|off-balance\s+sheet)'
            ])
            if contract_text:
                highlights['contractual_obligations'] = contract_text[:2000]
            
            data.mda_highlights = highlights
    
    async def _extract_risk_factors(self, text: str, data: AnnualReportData):
        """Extract risk factors."""
        risk_patterns = [
            r'(?i)(?:item\s+1a\.?\s*risk\s+factors|risk\s+factors)',
        ]
        
        risk_text = self._extract_section(text, risk_patterns, max_chars=30000)
        
        if risk_text:
            # Extract individual risk factors (often numbered or bulleted)
            risk_factors = re.findall(
                r'(?:^|\n)\s*(?:\d+\.|-\s+|\*\s+)([^\n]{50,500})',
                risk_text
            )
            data.risk_factors = [r.strip() for r in risk_factors[:50]]
    
    async def _extract_capital_allocation(self, text: str, data: AnnualReportData):
        """Extract capital allocation information."""
        # Dividends
        div_match = re.search(
            r'(?:dividend|dividends)\s+(?:of|was|is|declared|paid)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:per\s+share|per share)?',
            text, re.IGNORECASE
        )
        if div_match:
            data.dividends = float(div_match.group(1).replace(',', ''))
        
        # Share buybacks
        buyback_match = re.search(
            r'(?:repurchase|buyback|treasury\s+stock)\s+(?:of|was|is|amounted\s+to)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if buyback_match:
            value = float(buyback_match.group(1).replace(',', ''))
            if 'billion' in buyback_match.group(0).lower():
                value *= 1e9
            elif 'million' in buyback_match.group(0).lower():
                value *= 1e6
            data.buybacks = value
        
        # CapEx
        capex_match = re.search(
            r'(?:capital\s+expenditures?|capex)\s+(?:of|was|is|amounted\s+to)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if capex_match:
            value = float(capex_match.group(1).replace(',', ''))
            if 'billion' in capex_match.group(0).lower():
                value *= 1e9
            elif 'million' in capex_match.group(0).lower():
                value *= 1e6
            data.capex = value
    
    async def _extract_guidance(self, text: str, data: AnnualReportData):
        """Extract forward-looking guidance."""
        guidance_patterns = [
            r'(?i)(?:guidance|outlook|expects?|anticipates?|projects?|targets?)\s+(?:for\s+)?(?:20\d{2}|fiscal\s+20\d{2}|next\s+year|next\s+quarter)',
        ]
        
        guidance_text = self._extract_section(text, guidance_patterns, max_chars=10000)
        
        if guidance_text:
            # Extract specific guidance metrics
            guidance = {}
            
            # Revenue guidance
            rev_guidance = re.search(
                r'(?:revenue|sales)\s+(?:guidance|outlook|expected|expected\s+to\s+be)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
                guidance_text, re.IGNORECASE
            )
            if rev_guidance:
                guidance['revenue'] = rev_guidance.group(0)
            
            # EPS guidance
            eps_guidance = re.search(
                r'(?:eps|earnings\s+per\s+share)\s+(?:guidance|outlook|expected\s+to\s+be)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)',
                guidance_text, re.IGNORECASE
            )
            if eps_guidance:
                guidance['eps'] = eps_guidance.group(0)
            
            # Margin guidance
            margin_guidance = re.search(
                r'(?:gross|operating)\s+margin\s+(?:guidance|outlook|expected\s+to\s+be)\s*[:\\-]?\s*([\d,\.]+)\s*%',
                guidance_text, re.IGNORECASE
            )
            if margin_guidance:
                guidance[f'{margin_guidance.group(0).split()[0]}_margin'] = margin_guidance.group(0)
            
            data.guidance = guidance
    
    def _extract_section(self, text: str, patterns: list[str], max_chars: int = 5000) -> str:
        """Extract a section based on patterns."""
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                start = match.start()
                end = min(start + max_chars, len(text))
                return text[start:end]
        return ""
    
    def _extract_subsection(self, text: str, patterns: list[str]) -> str:
        """Extract a subsection within a larger section."""
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                start = match.start()
                end = min(start + 5000, len(text))
                return text[start:end]
        return ""
    
    def _calculate_confidence(self, data: AnnualReportData) -> float:
        """Calculate extraction confidence score."""
        score = 0.0
        total = 0
        
        # Check each field
        fields = [
            ('company_name', data.company_name),
            ('ticker', data.ticker),
            ('fiscal_year', data.fiscal_year),
            ('revenue', data.revenue),
            ('net_income', data.net_income),
            ('eps', data.eps),
            ('business_description', data.business_description),
            ('segments', data.segments),
            ('risk_factors', data.risk_factors),
            ('mda_highlights', data.mda_highlights),
        ]
        
        for name, value in fields:
            total += 1
            if value:
                score += 1
        
        return score / total if total > 0 else 0.0


class QuarterlyReportParser:
    """Parse quarterly reports (10-Q)."""
    
    async def parse(self, file_path: Path) -> QuarterlyReportData:
        result = await parse_financial_pdf(file_path)
        
        if not result.success:
            raise ValueError(f"Failed to parse {file_path}: {result.errors}")
        
        full_text = "\n\n".join(page.text for page in result.pages)
        
        data = QuarterlyReportData(
            company_name="",
            filing_date=None
        )
        
        # Extract quarter and year
        quarter_match = re.search(r'(?:q[1-4]|quarter\s+[1-4])\s+(?:20\d{2})', result.pages[0].text, re.IGNORECASE)
        if quarter_match:
            q_text = quarter_match.group(0).upper()
            q_match = re.search(r'Q([1-4])', q_text)
            y_match = re.search(r'(20\d{2})', q_text)
            if q_match:
                data.quarter = f"Q{q_match.group(1)}"
            if y_match:
                data.fiscal_year = int(y_match.group(1))
        
        # Extract financials
        text = "\n\n".join(page.text for page in result.pages)
        
        # Revenue
        rev_match = re.search(
            r'total\s+revenue\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            full_text, re.IGNORECASE
        )
        if rev_match:
            value = float(rev_match.group(1).replace(',', ''))
            if 'billion' in rev_match.group(0).lower():
                value *= 1e9
            elif 'million' in rev_match.group(0).lower():
                value *= 1e6
            data.revenue = value
        
        # Net income
        ni_match = re.search(
            r'net income\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            text, re.IGNORECASE
        )
        if ni_match:
            value = float(ni_match.group(1).replace(',', ''))
            if 'billion' in ni_match.group(0).lower():
                value *= 1e9
            elif 'million' in ni_match.group(0).lower():
                value *= 1e6
            data.net_income = value
        
        # EPS
        eps_match = re.search(
            r'(?:eps|earnings per share)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)',
            full_text, re.IGNORECASE
        )
        if eps_match:
            data.eps = float(eps_match.group(1).replace(',', ''))
        
        return data


class InvestorPresentationParser:
    """
    Parse investor presentations (PDF, PowerPoint).
    
    Extracts:
    - Slide content and structure
    - Key financial highlights
    - Strategic initiatives
    - Growth drivers
    - Guidance
    - Capital allocation
    - ESG highlights
    """
    
    def __init__(self):
        pass
    
    async def parse(self, file_path: Path) -> InvestorPresentationData:
        """Parse investor presentation."""
        result = await parse_financial_pdf(file_path)
        
        if not result.success:
            raise ValueError(f"Failed to parse {file_path}: {result.errors}")
        
        data = InvestorPresentationData(
            source_file=str(file_path),
            total_slides=len(result.pages)
        )
        
        # Extract text from all slides
        full_text = "\n\n".join(page.text for page in result.pages)
        
        # Extract metadata
        first_page = result.pages[0].text if result.pages else ""
        data.company_name = self._extract_company_name(first_page)
        
        # Try to find date
        date_match = re.search(r'(?:january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{1,2},?\s+20\d{2}', first_page, re.IGNORECASE)
        if date_match:
            try:
                data.presentation_date = datetime.strptime(date_match.group(0), "%B %d, %Y").date()
            except:
                pass
        
        # Extract title
        lines = first_page.split('\n')[:10]
        for line in lines:
            if len(line.strip()) > 10 and not line.strip().isdigit():
                data.presentation_title = line.strip()
                break
        
        # Process each slide
        for i, page in enumerate(result.pages):
            slide_data = await self._process_slide(page, i + 1)
            data.slides.append(slide_data)
        
        # Extract key highlights
        await self._extract_key_highlights(data, full_text)
        
        # Extract financial highlights
        await self._extract_financial_highlights(full_text, data)
        
        # Extract strategic initiatives
        await self._extract_strategic_initiatives(data, full_text)
        
        # Extract guidance
        await self._extract_guidance(full_text, data)
        
        # Extract capital allocation
        await self._extract_capital_allocation(data, full_text)
        
        return data
    
    def _process_slide(self, page, slide_number: int) -> dict:
        """Process a single slide."""
        # Detect if slide is a chart/table vs text
        has_tables = len(page.tables) > 0
        text_length = len(page.text)
        
        slide_type = "text"
        if has_tables:
            slide_type = "table"
        elif text_length < 100:
            slide_type = "visual"
        
        return {
            "slide_number": slide_number,
            "type": slide_type,
            "text": page.text[:2000],
            "tables_count": len(page.tables),
            "tables": [
                {
                    "index": t.table_index,
                    "headers": t.headers,
                    "rows": t.rows[:10] if t.rows else []
                }
                for t in page.tables
            ] if page.tables else []
        }
    
    def _extract_company_name(self, text: str) -> str:
        """Extract company name from first page."""
        lines = text.split('\n')[:20]
        for line in lines:
            if any(keyword in line.lower() for keyword in ['corporation', 'inc.', 'inc', 'company', 'ltd', 'llc', 'corp']):
                return line.strip()
        return ""
    
    async def _extract_key_highlights(self, data: InvestorPresentationData, full_text: str):
        """Extract key highlights from presentation."""
        # Look for highlight sections
        highlight_patterns = [
            r'(?i)(?:key\s+highlights?|highlights?|key\s+points?|important\s+points?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:key\s+takeaways?|takeaways?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        for pattern in highlight_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                highlights_text = match.group(1)
                # Split into bullet points
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', highlights_text)
                data.key_highlights = [b.strip() for b in bullets[:10]]
                break
        
        # If no explicit highlights, look for bold/emphasized text
        if not data.key_highlights:
            bold_matches = re.findall(r'\*\*([^*]{20,200})\*\*', full_text)
            data.key_highlights = bold_matches[:10]
    
    async def _extract_financial_highlights(self, full_text: str, data: InvestorPresentationData):
        """Extract financial metrics shown in presentation."""
        # Common financial metrics in presentations
        metrics = {
            'revenue': r'(?:revenue|sales)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'net_income': r'net income\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'eps': r'(?:eps|earnings per share)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)',
            'gross_margin': r'gross margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%',
            'operating_margin': r'operating margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%',
            'free_cash_flow': r'free cash flow\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'ebitda': r'ebitda\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'cash': r'(?:cash\s+(?:and\s+)?(?:equivalents|and\s+short-term\s+investments))\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
        }
        
        for metric, pattern in metrics.items():
            match = re.search(pattern, full_text, re.IGNORECASE)
            if match:
                data.key_metrics[metric] = match.group(0)
    
    async def _extract_strategic_initiatives(self, data: InvestorPresentationData, full_text: str):
        """Extract strategic initiatives from presentation."""
        strategy_patterns = [
            r'(?i)(?:strategic\s+initiatives?|key\s+initiatives?|strategic\s+priorities?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:growth\s+drivers?|growth\s+strategies?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:strategic\s+focus|focus\s+areas?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        initiatives = []
        for pattern in strategy_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                # Split into bullet points
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', text)
                initiatives.extend([b.strip() for b in bullets[:5]])
        
        data.strategic_initiatives = initiatives[:10]
    
    async def _extract_growth_drivers(self, data: InvestorPresentationData, full_text: str):
        """Extract growth drivers mentioned."""
        driver_patterns = [
            r'(?i)(?:growth\s+drivers?|key\s+drivers?|drivers?\s+of\s+growth)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        for pattern in driver_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', text)
                data.growth_drivers = [b.strip() for b in bullets[:10]]
                break
    
    async def _extract_guidance(self, full_text: str, data: InvestorPresentationData):
        """Extract guidance from presentation."""
        guidance_patterns = [
            r'(?i)(?:guidance|outlook|targets?|goals?)\s+(?:for\s+)?(?:20\d{2}|fiscal\s+20\d{2}|next\s+year|next\s+quarter)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        guidance = {}
        for pattern in guidance_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                # Extract specific metrics
                for metric in ['revenue', 'eps', 'earnings', 'margin', 'capex', 'free cash flow', 'cash flow']:
                    match = re.search(
                        rf'{metric}\s+(?:of|target|expected|guidance)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?\s*%',
                        text, re.IGNORECASE
                    )
                    if match:
                        guidance[metric] = match.group(0)
                break
        
        data.guidance = guidance
    
    async def _extract_capital_allocation(self, data: InvestorPresentationData, full_text: str):
        """Extract capital allocation information."""
        allocation = {}
        
        # Dividends
        div_match = re.search(
            r'(?:dividend|dividends)\s+(?:of|yield|increase|raised)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:per\s+share|%)?',
            full_text, re.IGNORECASE
        )
        if div_match:
            allocation['dividend'] = div_match.group(0)
        
        # Buybacks
        bb_match = re.search(
            r'(?:repurchase|buyback|share\s+repurchase)\s+(?:of|amounting\s+to|authorized)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            full_text, re.IGNORECASE
        )
        if bb_match:
            allocation['buyback'] = bb_match.group(0)
        
        # CapEx
        capex_match = re.search(
            r'(?:capital\s+expenditures?|capex)\s+(?:of|expected|projected|target)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            full_text, re.IGNORECASE
        )
        if capex_match:
            allocation['capex'] = capex_match.group(0)
        
        data.capital_allocation = allocation
    
    async def parse_pptx(self, file_path: Path) -> InvestorPresentationData:
        """Parse PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not available, falling back to PDF parsing")
            return await self.parse(file_path)
        
        prs = Presentation(file_path)
        
        data = InvestorPresentationData(
            source_file=str(file_path),
            total_slides=len(prs.slides)
        )
        
        full_text = ""
        
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text + "\n"
            
            full_text += slide_text + "\n\n"
            
            data.slides.append({
                "slide_number": i + 1,
                "text": slide_text[:2000],
                "shapes_count": len(slide.shapes)
            })
        
        # Parse using text-based methods
        return await self.parse_text(full_text, data.company_name)
    
    async def parse_text(self, full_text: str, company_name: str = "") -> InvestorPresentationData:
        """Parse from raw text."""
        data = InvestorPresentationData(
            company_name=company_name,
            total_slides=0
        )
        
        await self._extract_key_highlights(data, full_text)
        await self._extract_financial_highlights(full_text, data)
        await self._extract_strategic_initiatives(data, full_text)
        await self._extract_growth_drivers(data, full_text)
        await self._extract_guidance(full_text, data)
        await self._extract_capital_allocation(data, full_text)
        
        return data


class FinancialDocumentParser:
    """
    Main parser orchestrating all financial document types.
    
    Automatically detects document type and routes to appropriate parser.
    """
    
    def __init__(self):
        self.annual_report_parser = AnnualReportParser()
        self.quarterly_parser = QuarterlyReportParser()
        self.presentation_parser = InvestorPresentationParser()
        self.transcript_parser = EarningsCallTranscriptParser()
        
        # Document type detection patterns
        self.type_patterns = {
            '10k': [
                r'form\s+10-?k',
                r'annual\s+report\s+(?:on\s+form\s+)?10-?k',
                r'10-?k',
            ],
            '10q': [
                r'form\s+10-?q',
                r'quarterly\s+report\s+(?:on\s+form\s+)?10-?q',
                r'10-?q',
            ],
            'annual_report': [
                r'annual\s+report',
                r'year\s+ended',
                r'fiscal\s+year\s+20\d{2}',
            ],
            '10q': [
                r'quarterly\s+report',
                r'quarter\s+ended',
                r'form\s+10-?q',
            ],
            'earnings_transcript': [
                r'earnings\s+call',
                r'q[1-4]\s+earnings',
                r'conference\s+call',
                r'q&a',
            ],
            'investor_presentation': [
                r'investor\s+presentation',
                r'analyst\s+(?:day|presentation)',
                r'investor\s+day',
                r'presentation\s+to\s+analysts',
            ],
            'earnings_presentation': [
                r'earnings\s+presentation',
                r'quarterly\s+earnings\s+presentation',
            ],
        }
    
    def detect_document_type(self, file_path: Path, text: str = "") -> str:
        """Detect document type from filename and content."""
        filename = file_path.name.lower()
        text_lower = text.lower()[:10000]
        
        for doc_type, patterns in self.type_patterns.items():
            for pattern in patterns:
                if re.search(pattern, filename, re.IGNORECASE) or re.search(pattern, text_lower):
                    return doc_type
        
        return "financial_document"
    
    async def parse(self, file_path: Path) -> dict:
        """
        Parse any financial document and return structured data.
        
        Returns a dict with:
        - document_type
        - data: parsed structured data
        - raw_result: ParserResult
        """
        # Parse with financial PDF parser
        result = await parse_financial_pdf(file_path)
        
        if not result.success:
            return {
                "document_type": "unknown",
                "data": None,
                "raw_result": result,
                "error": result.errors
            }
        
        # Extract text
        full_text = "\n\n".join(page.text for page in result.pages)
        
        # Detect document type
        doc_type = self.detect_document_type(file_path, full_text)
        
        # Parse based on type
        if doc_type in ['10k', 'annual_report']:
            data = await self.annual_report_parser.parse(file_path)
            return {"document_type": "annual_report", "data": data, "raw_result": result}
        
        elif doc_type in ['10q', 'quarterly_report']:
            data = await self.quarterly_parser.parse(file_path)
            return {"document_type": "quarterly_report", "data": data, "raw_result": result}
        
        elif doc_type == 'earnings_transcript':
            data = await self.transcript_parser.parse(file_path)
            return {"document_type": "earnings_transcript", "data": data, "raw_result": result}
        
        elif doc_type in ['investor_presentation', 'earnings_presentation']:
            data = await self.presentation_parser.parse(file_path)
            return {"document_type": "investor_presentation", "data": data, "raw_result": result}
        
        # Default: return raw parsed result
        return {
            "document_type": "financial_document",
            "data": {
                "pages": len(result.pages),
                "full_text": "\n\n".join(p.text for p in result.pages)
            },
            "raw_result": result
        }
    
    async def parse_batch(self, file_paths: list[Path]) -> list[dict]:
        """Parse multiple documents."""
        results = []
        for path in file_paths:
            try:
                result = await self.parse(path)
                results.append(result)
            except Exception as e:
                logger.error(f"Failed to parse {path}: {e}")
                results.append({"document_type": "error", "error": str(e), "file": str(path)})
        return results


# Factory
class FinancialDocumentParserFactory:
    """Factory for creating financial document parsers."""
    
    @staticmethod
    def create_parser(parser_type: str):
        """Create parser by type."""
        parsers = {
            'annual_report': AnnualReportParser(),
            'quarterly_report': QuarterlyReportParser(),
            'investor_presentation': InvestorPresentationParser(),
            'earnings_transcript': EarningsCallTranscriptParser(),
            'general': FinancialDocumentParser(),
        }
        
        return parsers.get(parser_type, FinancialDocumentParser())
    
    @staticmethod
    def create_all():
        """Create all parsers."""
        return {
            'annual_report': AnnualReportParser(),
            'quarterly_report': QuarterlyReportParser(),
            'investor_presentation': InvestorPresentationParser(),
            'earnings_transcript': EarningsCallTranscriptParser(),
            'financial_document': FinancialDocumentParser(),
        }


# Export
__all__ = [
    "AnnualReportData",
    "QuarterlyReportData",
    "InvestorPresentationData",
    "AnnualReportParser",
    "QuarterlyReportParser",
    "InvestorPresentationParser",
    "FinancialDocumentParser",
    "FinancialDocumentParserFactory",
    "EarningsCallTranscriptParser",
    "EarningsCallProcessor",
    "EarningsCallTranscript",
    "Speaker",
    "TranscriptSection",
    "QAExchange",
    "GuidanceItem",
    "KeyMetric",
]