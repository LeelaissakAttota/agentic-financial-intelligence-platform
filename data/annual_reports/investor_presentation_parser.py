"""
Investor Presentation Parser.

Parses investor presentations (PDF, PowerPoint) to extract:
- Slide content and structure
- Key financial highlights
- Strategic initiatives
- Growth drivers
- Guidance
- Capital allocation
- ESG highlights
"""

import asyncio
import logging
import re
from dataclasses import dataclass, field
from datetime import date
from pathlib import Path
from typing import Any, Optional

from data.financial_documents.parser import parse_financial_pdf, ParserResult

logger = logging.getLogger(__name__)


@dataclass
class InvestorPresentationData:
    """Structured data extracted from investor presentations."""
    company_name: str
    ticker: Optional[str] = None
    presentation_date: Optional[date] = None
    presentation_title: str = ""
    
    # Slides data
    slides: list[dict] = field(default_factory=list)
    
    # Key sections
    key_highlights: list[str] = field(default_factory=list)
    financial_highlights: dict[str, Any] = field(default_factory=dict)
    strategic_initiatives: list[str] = field(default_factory=list)
    growth_drivers: list[str] = field(default_factory=list)
    risks: list[str] = field(default_factory=list)
    
    # Financial metrics shown
    key_metrics: dict[str, Any] = field(default_factory=dict)
    
    # Guidance
    guidance: dict[str, Any] = field(default_factory=dict)
    
    # Capital allocation
    capital_allocation: dict[str, Any] = field(default_factory=dict)
    
    # ESG
    esg_highlights: dict[str, Any] = field(default_factory=dict)
    
    # Source
    source_file: Optional[str] = None
    total_slides: int = 0


class InvestorPresentationParser:
    """
    Parse investor presentations (PDF, PowerPoint).
    
    Extracts:
    - Slide content and structure
    - Key financial highlights
    - Strategic initiatives
    - Growth drivers
    - Guidance
    - Capital allocation
    - ESG highlights
    """
    
    def __init__(self):
        pass
    
    async def parse(self, file_path: Path) -> InvestorPresentationData:
        """Parse investor presentation."""
        result = await parse_financial_pdf(file_path)
        
        if not result.success:
            raise ValueError(f"Failed to parse {file_path}: {result.errors}")
        
        data = InvestorPresentationData(
            source_file=str(file_path),
            total_slides=len(result.pages)
        )
        
        # Extract text from all slides
        full_text = "\n\n".join(page.text for page in result.pages)
        
        # Extract metadata
        first_page = result.pages[0].text if result.pages else ""
        data.company_name = self._extract_company_name(first_page)
        
        # Try to find date
        date_match = re.search(r'(?:january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{1,2},?\s+20\d{2}', first_page, re.IGNORECASE)
        if date_match:
            try:
                data.presentation_date = datetime.strptime(date_match.group(0), "%B %d, %Y").date()
            except:
                pass
        
        # Extract title
        lines = first_page.split('\n')[:10]
        for line in lines:
            if len(line.strip()) > 10 and not line.strip().isdigit():
                data.presentation_title = line.strip()
                break
        
        # Process each slide
        for i, page in enumerate(result.pages):
            slide_data = await self._process_slide(page, i + 1)
            data.slides.append(slide_data)
        
        # Extract key highlights
        await self._extract_key_highlights(data, full_text)
        
        # Extract financial highlights
        await self._extract_financial_highlights(full_text, data)
        
        # Extract strategic initiatives
        await self._extract_strategic_initiatives(data, full_text)
        
        # Extract guidance
        await self._extract_guidance(full_text, data)
        
        # Extract capital allocation
        await self._extract_capital_allocation(data, full_text)
        
        return data
    
    def _process_slide(self, page, slide_number: int) -> dict:
        """Process a single slide."""
        # Detect if slide is a chart/table vs text
        has_tables = len(page.tables) > 0
        text_length = len(page.text)
        
        slide_type = "text"
        if has_tables:
            slide_type = "table"
        elif text_length < 100:
            slide_type = "visual"
        
        return {
            "slide_number": slide_number,
            "type": slide_type,
            "text": page.text[:2000],
            "tables_count": len(page.tables),
            "tables": [
                {
                    "index": t.table_index,
                    "headers": t.headers,
                    "rows": t.rows[:10] if t.rows else []
                }
                for t in page.tables
            ] if page.tables else []
        }
    
    def _extract_company_name(self, text: str) -> str:
        """Extract company name from first page."""
        lines = text.split('\n')[:20]
        for line in lines:
            if any(keyword in line.lower() for keyword in ['corporation', 'inc.', 'inc', 'company', 'ltd', 'llc', 'corp']):
                return line.strip()
        return ""
    
    async def _extract_key_highlights(self, data: InvestorPresentationData, full_text: str):
        """Extract key highlights from presentation."""
        # Look for highlight sections
        highlight_patterns = [
            r'(?i)(?:key\s+highlights?|highlights?|key\s+points?|important\s+points?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:key\s+takeaways?|takeaways?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        for pattern in highlight_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                highlights_text = match.group(1)
                # Split into bullet points
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', highlights_text)
                data.key_highlights = [b.strip() for b in bullets[:10]]
                break
        
        # If no explicit highlights, look for bold/emphasized text
        if not data.key_highlights:
            bold_matches = re.findall(r'\*\*([^*]{20,200})\*\*', full_text)
            data.key_highlights = bold_matches[:10]
    
    async def _extract_financial_highlights(self, full_text: str, data: InvestorPresentationData):
        """Extract financial metrics shown in presentation."""
        # Common financial metrics in presentations
        metrics = {
            'revenue': r'(?:revenue|sales)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'net_income': r'net income\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'eps': r'(?:eps|earnings per share)\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)',
            'gross_margin': r'gross margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%',
            'operating_margin': r'operating margin\s+(?:of|was|is)\s*([\d,\.]+)\s*%',
            'free_cash_flow': r'free cash flow\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'ebitda': r'ebitda\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            'cash': r'(?:cash\s+(?:and\s+)?(?:equivalents|and\s+short-term\s+investments))\s+(?:of|was|is)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
        }
        
        for metric, pattern in metrics.items():
            match = re.search(pattern, full_text, re.IGNORECASE)
            if match:
                data.key_metrics[metric] = match.group(0)
    
    async def _extract_strategic_initiatives(self, data: InvestorPresentationData, full_text: str):
        """Extract strategic initiatives from presentation."""
        strategy_patterns = [
            r'(?i)(?:strategic\s+initiatives?|key\s+initiatives?|strategic\s+priorities?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:growth\s+drivers?|growth\s+strategies?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
            r'(?i)(?:strategic\s+focus|focus\s+areas?)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        initiatives = []
        for pattern in strategy_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                # Split into bullet points
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', text)
                initiatives.extend([b.strip() for b in bullets[:5]])
        
        data.strategic_initiatives = initiatives[:10]
    
    async def _extract_growth_drivers(self, data: InvestorPresentationData, full_text: str):
        """Extract growth drivers mentioned."""
        driver_patterns = [
            r'(?i)(?:growth\s+drivers?|key\s+drivers?|drivers?\s+of\s+growth)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        for pattern in driver_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                bullets = re.findall(r'(?:^|\n)[\-\*\•]\s*([^\n]{10,200})', text)
                data.growth_drivers = [b.strip() for b in bullets[:10]]
                break
    
    async def _extract_guidance(self, full_text: str, data: InvestorPresentationData):
        """Extract guidance from presentation."""
        guidance_patterns = [
            r'(?i)(?:guidance|outlook|targets?|goals?)\s+(?:for\s+)?(?:20\d{2}|fiscal\s+20\d{2}|next\s+year|next\s+quarter)[:\\s]*\n?(.*?)(?:\n\n|\Z)',
        ]
        
        guidance = {}
        for pattern in guidance_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
            if match:
                text = match.group(1)
                # Extract specific metrics
                for metric in ['revenue', 'eps', 'earnings', 'margin', 'capex', 'free cash flow', 'cash flow']:
                    match = re.search(
                        rf'{metric}\s+(?:of|target|expected|guidance)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?\s*%',
                        text, re.IGNORECASE
                    )
                    if match:
                        guidance[metric] = match.group(0)
                break
        
        data.guidance = guidance
    
    async def _extract_capital_allocation(self, data: InvestorPresentationData, full_text: str):
        """Extract capital allocation information."""
        allocation = {}
        
        # Dividends
        div_match = re.search(
            r'(?:dividend|dividends)\s+(?:of|yield|increase|raised)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:per\s+share|%)?',
            full_text, re.IGNORECASE
        )
        if div_match:
            allocation['dividend'] = div_match.group(0)
        
        # Buybacks
        bb_match = re.search(
            r'(?:repurchase|buyback|share\s+repurchase)\s+(?:of|amounting\s+to|authorized)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            full_text, re.IGNORECASE
        )
        if bb_match:
            allocation['buyback'] = bb_match.group(0)
        
        # CapEx
        capex_match = re.search(
            r'(?:capital\s+expenditures?|capex)\s+(?:of|expected|projected|target)\s*[:\\-]?\s*[\$£€]?\s*([\d,\.]+)\s*(?:million|billion|b|m|mm)?',
            full_text, re.IGNORECASE
        )
        if capex_match:
            allocation['capex'] = capex_match.group(0)
        
        data.capital_allocation = allocation
    
    async def parse_pptx(self, file_path: Path) -> InvestorPresentationData:
        """Parse PowerPoint file using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not available, falling back to PDF parsing")
            return await self.parse(file_path)
        
        prs = Presentation(file_path)
        
        data = InvestorPresentationData(
            source_file=str(file_path),
            total_slides=len(prs.slides)
        )
        
        full_text = ""
        
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text + "\n"
            
            full_text += slide_text + "\n\n"
            
            data.slides.append({
                "slide_number": i + 1,
                "text": slide_text[:2000],
                "shapes_count": len(slide.shapes)
            })
        
        # Parse using text-based methods
        return await self.parse_text(full_text, data.company_name)
    
    async def parse_text(self, full_text: str, company_name: str = "") -> InvestorPresentationData:
        """Parse from raw text."""
        data = InvestorPresentationData(
            company_name=company_name,
            total_slides=0
        )
        
        await self._extract_key_highlights(data, full_text)
        await self._extract_financial_highlights(full_text, data)
        await self._extract_strategic_initiatives(data, full_text)
        await self._extract_growth_drivers(data, full_text)
        await self._extract_guidance(full_text, data)
        await self._extract_capital_allocation(data, full_text)
        
        return data


# Export
__all__ = [
    "InvestorPresentationData",
    "InvestorPresentationParser",
]